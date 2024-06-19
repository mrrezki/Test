from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Define a function to add a slide with a title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Use layout with title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Title Slide
slide_layout = prs.slide_layouts[0]  # Use title slide layout
slide = prs.slides.add_slide(slide_layout)
title_placeholder = slide.shapes.title
subtitle_placeholder = slide.placeholders[1]
title_placeholder.text = "Typical Use Cases for Generative AI"
subtitle_placeholder.text = "Exploring Applications Across Various Industries"

# Slide 2: Content Creation
add_slide(
    "Content Creation",
    "Text Generation: Automated article writing, blog posts, and news reports.\n"
    "Creative Writing: Generating poetry, short stories, and scripts.\n"
    "Summarization: Creating concise summaries of long documents."
)

# Slide 3: Image and Video Generation
add_slide(
    "Image and Video Generation",
    "Art and Design: Creating original artworks, graphic designs, and digital content.\n"
    "Video Production: Generating realistic videos and animations from textual descriptions.\n"
    "Photo Editing: Enhancing or transforming images (e.g., colorizing black and white photos)."
)

# Slide 4: Chatbots and Virtual Assistants
add_slide(
    "Chatbots and Virtual Assistants",
    "Customer Support: Providing automated responses to customer inquiries.\n"
    "Personal Assistants: Managing schedules, sending reminders, and handling basic tasks."
)

# Slide 5: Gaming
add_slide(
    "Gaming",
    "Procedural Content Generation: Creating new levels, characters, and environments.\n"
    "Storytelling: Generating dynamic narratives based on player actions."
)

# Slide 6: Healthcare
add_slide(
    "Healthcare",
    "Drug Discovery: Generating new molecular structures for potential medications.\n"
    "Medical Imaging: Enhancing or generating diagnostic images (e.g., MRI, CT scans).\n"
    "Personalized Treatment Plans: Analyzing patient data to recommend treatments."
)

# Slide 7: Finance
add_slide(
    "Finance",
    "Fraud Detection: Identifying unusual patterns that may indicate fraudulent activity.\n"
    "Algorithmic Trading: Generating trading strategies based on historical data.\n"
    "Risk Management: Predicting and managing financial risks."
)

# Slide 8: Education
add_slide(
    "Education",
    "Personalized Learning: Creating customized educational materials for students.\n"
    "Automated Tutoring: Providing explanations and assistance on various subjects."
)

# Slide 9: Marketing and Advertising
add_slide(
    "Marketing and Advertising",
    "Ad Copy Generation: Creating engaging and targeted advertising content.\n"
    "Market Analysis: Generating insights and reports from large datasets."
)

# Slide 10: Music and Audio
add_slide(
    "Music and Audio",
    "Music Composition: Creating original music tracks in various genres.\n"
    "Sound Effects: Generating realistic sound effects for movies and games.\n"
    "Voice Synthesis: Producing natural-sounding speech for voiceovers and assistants."
)

# Slide 11: Software Development
add_slide(
    "Software Development",
    "Code Generation: Automating parts of the coding process by generating code snippets.\n"
    "Automated Testing: Generating test cases and scenarios for software applications."
)

# Slide 12: Personalization
add_slide(
    "Personalization",
    "Recommendations: Providing personalized product, content, or service recommendations.\n"
    "Customization: Generating personalized experiences for users based on their preferences."
)

# Slide 13: Scientific Research
add_slide(
    "Scientific Research",
    "Data Analysis: Generating hypotheses and insights from complex datasets.\n"
    "Simulation: Creating simulations to model scientific phenomena and predict outcomes."
)

# Save the presentation
prs.save("Generative_AI_Use_Cases_Presentation.pptx")
